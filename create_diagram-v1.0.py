### python-pptx 1.0.2 环境中稳定运行的、具备智能连接点选择的脚本
import json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR 

# --- 导入 Enum 常量 (保持兼容性回退) ---
try:
    from pptx.enum.connector import MSO_CONNECTOR_TYPE 
    from pptx.enum.line import MSO_ARROWHEAD_STYLE 

    CONNECTOR_STRAIGHT = MSO_CONNECTOR_TYPE.STRAIGHT
    ARROWHEAD_TRIANGLE = MSO_ARROWHEAD_STYLE.TRIANGLE
    ARROWHEAD_NONE = MSO_ARROWHEAD_STYLE.NONE
    
except ImportError:
    # 警告：导入失败。回退到硬编码常量以保持 Enum 语义。
    print(f"警告：导入必要的Enum失败 (No module named 'pptx.enum.connector')。正在回退到硬编码常量以继续执行。")
    class MSO_CONNECTOR_TYPE:
        STRAIGHT = 1
    class MSO_ARROWHEAD_STYLE:
        NONE = 0
        TRIANGLE = 2 
    
    CONNECTOR_STRAIGHT = MSO_CONNECTOR_TYPE.STRAIGHT
    ARROWHEAD_TRIANGLE = MSO_ARROWHEAD_STYLE.TRIANGLE
    ARROWHEAD_NONE = MSO_ARROWHEAD_STYLE.NONE


# --- 配置 ---
CANVAS_WIDTH_PX = 1000
CANVAS_HEIGHT_PX = 800
SLIDE_WIDTH_INCHES = 10
SLIDE_HEIGHT_INCHES = 7.5
TOLERANCE_INCHES = Inches(0.5) # 0.5英寸的对齐容错范围

# 🚨 最终修正：使用确定的有效连接点索引 [0, 1, 2, 3]
CONNECT_LEFT = 0 
CONNECT_TOP = 1
CONNECT_RIGHT = 2
CONNECT_BOTTOM = 3


# 像素值到英寸的转换函数 (保持不变)
def px_to_inches(px_value, axis_size_px):
    if axis_size_px == CANVAS_WIDTH_PX:
        return Inches(px_value * SLIDE_WIDTH_INCHES / CANVAS_WIDTH_PX)
    else:
        return Inches(px_value * SLIDE_HEIGHT_INCHES / CANVAS_HEIGHT_PX)

# RGB字符串转RGBColor对象 (保持不变)
def rgb_string_to_pptx_color(rgb_str):
    try:
        r, g, b = map(int, rgb_str.strip('RGB()').split(','))
        return RGBColor(r, g, b)
    except:
        return RGBColor(0, 0, 0)

# 从JSON文件中加载数据 (保持不变)
try:
    with open('diagram_data.json', 'r', encoding='utf-8') as f:
        data = json.load(f)
except FileNotFoundError:
    print("错误：未找到 'diagram_data.json' 文件。请确保JSON文件已保存。")
    exit()
except json.JSONDecodeError:
    print("错误：'diagram_data.json' 文件格式不正确。")
    exit()

# --- 核心优化函数 ---
def get_optimal_connection_points(start_shape, end_shape):
    """
    根据形状的相对位置和对齐容错，智能确定连接器的最佳连接点。
    使用确定的 [0, 1, 2, 3] 索引。
    """
    start_center_x = start_shape.left + start_shape.width / 2
    start_center_y = start_shape.top + start_shape.height / 2
    
    end_center_x = end_shape.left + end_shape.width / 2
    end_center_y = end_shape.top + end_shape.height / 2
    
    dx = abs(start_center_x - end_center_x)
    
    # 检查X轴对齐：如果X轴中心点偏差小于容错范围，视为垂直关系
    if dx < TOLERANCE_INCHES:
        # 垂直关系更突出，强制使用顶部/底部
        if start_center_y < end_center_y:
            # S在上，E在下（最常见）
            return CONNECT_BOTTOM, CONNECT_TOP 
        else:
            # S在下，E在上 (回流)
            return CONNECT_TOP, CONNECT_BOTTOM
    else:
        # X轴偏差较大，视为水平关系，使用左侧/右侧
        if start_center_x < end_center_x:
            # S在左，E在右
            return CONNECT_RIGHT, CONNECT_LEFT
        else:
            # S在右，E在左 (回流)
            return CONNECT_LEFT, CONNECT_RIGHT

def create_powerpoint_diagram(data):
    """根据JSON数据创建PowerPoint幻灯片"""
    
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_WIDTH_INCHES)
    prs.slide_height = Inches(SLIDE_HEIGHT_INCHES)
    
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    
    shapes_map = {}
    
    # 1. 创建所有图形元素 (保持不变)
    for element in data['elements']:
        text = element['text']
        
        if 'position' not in element or 'dimensions' not in element:
            continue
            
        try:
            x_px_str, y_px_str = element['position'].strip('[]').split(',')
            w_px_str, h_px_str = element['dimensions'].strip('[]').split(',')
            
            x_px, y_px = int(x_px_str), int(y_px_str)
            w_px, h_px = int(w_px_str), int(h_px_str)
            
            left = px_to_inches(x_px, CANVAS_WIDTH_PX)
            top = px_to_inches(y_px, CANVAS_HEIGHT_PX)
            width = px_to_inches(w_px, CANVAS_WIDTH_PX)
            height = px_to_inches(h_px, CANVAS_HEIGHT_PX)
        except:
            if 'id' not in element: 
                continue 
            continue

        shape_type = MSO_AUTO_SHAPE_TYPE.RECTANGLE
        if element['type'].endswith('rounded_rectangle'):
            shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE
        elif element['type'].startswith('circle'):
            shape_type = MSO_AUTO_SHAPE_TYPE.OVAL
        
        if element['type'].endswith('text') or 'independent_text' in element['type']:
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.text = text
            p = tf.paragraphs[0]
            p.font.size = Pt(9)
            continue
            
        shape = slide.shapes.add_shape(shape_type, left, top, width, height)

        fill = shape.fill
        fill.solid()
        fill.fore_color.rgb = rgb_string_to_pptx_color(element.get('color', 'RGB(200, 200, 200)'))
        line = shape.line
        line.color.rgb = RGBColor(0, 0, 0)
        line.width = Pt(0.75)

        text_frame = shape.text_frame
        text_frame.text = text
        text_frame.word_wrap = True
        text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE 
        text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        
        p = text_frame.paragraphs[0]
        p.font.size = Pt(10)
        p.font.color.rgb = RGBColor(0, 0, 0)
        p.alignment = MSO_ANCHOR.TOP
        
        if 'id' in element:
            shapes_map[element['id']] = shape
            
    # 2. 创建连接线 (箭头和无箭头的线段)
    for relationship in data.get('relationships', []):
        
        from_id = relationship.get('from')
        to_id = relationship.get('to')
        link_type = relationship.get('type')
        
        if from_id in shapes_map and to_id in shapes_map:
            start_shape = shapes_map[from_id]
            end_shape = shapes_map[to_id]

            if link_type in ['arrow', 'line', 'arrow_flow_down']:
                # --- 紧密连接 ---
                has_arrow = 'arrow' in link_type
                
                connector = slide.shapes.add_connector(
                    CONNECTOR_STRAIGHT, Inches(0), Inches(0), Inches(0), Inches(0)
                )

                # 核心优化点：使用精确的连接点逻辑 [0, 1, 2, 3]
                start_conn_pt, end_conn_pt = get_optimal_connection_points(start_shape, end_shape)
                
                connector.begin_connect(start_shape, start_conn_pt)
                connector.end_connect(end_shape, end_conn_pt)
                
                line = connector.line
                line.color.rgb = RGBColor(0, 0, 0)
                line.width = Pt(1.5)
                
                if has_arrow:
                    line.end_arrowhead_style = ARROWHEAD_TRIANGLE 
                else:
                    line.end_arrowhead_style = ARROWHEAD_NONE

            # 松散连接（疑问元素）
            elif 'loose_line' in relationship.get('link_type', ''):
                # 松散连接：S右侧连接到 E左侧 (使用 0 和 2)
                connector = slide.shapes.add_connector(
                    CONNECTOR_STRAIGHT, Inches(0), Inches(0), Inches(0), Inches(0)
                )
                
                connector.begin_connect(start_shape, CONNECT_RIGHT) # 2
                connector.end_connect(end_shape, CONNECT_LEFT)    # 0
                
                line = connector.line
                line.color.rgb = RGBColor(128, 128, 128) # 灰色
                line.width = Pt(1)
                line.end_arrowhead_style = ARROWHEAD_NONE


    output_path = '实例映射研讨会结果_最终完美版.pptx'
    prs.save(output_path)
    print(f"\n✅ 成功生成PowerPoint文件: {output_path}")

# 执行函数
create_powerpoint_diagram(data)