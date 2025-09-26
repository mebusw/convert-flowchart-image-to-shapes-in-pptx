import json
import math
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR 

# --- 兼容性导入及连接器类型定义 ---
try:
    from pptx.enum.connector import MSO_CONNECTOR_TYPE 
    from pptx.enum.line import MSO_ARROWHEAD_STYLE 

    CONNECTOR_ELBOW = MSO_CONNECTOR_TYPE.ELBOW 
    CONNECTOR_STRAIGHT = MSO_CONNECTOR_TYPE.STRAIGHT
    ARROWHEAD_TRIANGLE = MSO_ARROWHEAD_STYLE.TRIANGLE
    ARROWHEAD_NONE = MSO_ARROWHEAD_STYLE.NONE
    
except ImportError:
    print("警告：部分Enum导入失败，使用兼容性常量")
    class MSO_CONNECTOR_TYPE:
        STRAIGHT = 1
        ELBOW = 3
    class MSO_ARROWHEAD_STYLE:
        NONE = 0
        TRIANGLE = 2 
    
    CONNECTOR_ELBOW = MSO_CONNECTOR_TYPE.ELBOW 
    CONNECTOR_STRAIGHT = MSO_CONNECTOR_TYPE.STRAIGHT
    ARROWHEAD_TRIANGLE = MSO_ARROWHEAD_STYLE.TRIANGLE
    ARROWHEAD_NONE = MSO_ARROWHEAD_STYLE.NONE

# --- 配置参数 ---
CANVAS_WIDTH_PX = 1000
CANVAS_HEIGHT_PX = 800
SLIDE_WIDTH_INCHES = 10
SLIDE_HEIGHT_INCHES = 7.5

# 连接点索引常量
CONNECT_LEFT = 0    # 左侧连接点
CONNECT_TOP = 1     # 顶部连接点
CONNECT_RIGHT = 2   # 右侧连接点
CONNECT_BOTTOM = 3  # 底部连接点

# 优化的容错参数 (EMU 值用于精确计算，Inches 用于定义)
X_ALIGNMENT_TOLERANCE_IN = Inches(0.8)  # X轴对齐容错范围
Y_ALIGNMENT_TOLERANCE_IN = Inches(0.6)  # Y轴对齐容错范围
MIN_DISTANCE_THRESHOLD_IN = Inches(0.3) # 最小距离阈值

# ----------------------------------------------------
# --- 辅助函数定义（必须在主函数之前） ---
# ----------------------------------------------------

def get_emu_value(pptx_unit):
    """
    获取 pptx 单位对象（如 Inches, Pt）的原始 EMU 数值。
    用于确保比较时单位一致。
    """
    if hasattr(pptx_unit, '_emu'):
        return pptx_unit._emu
    return pptx_unit 


def px_to_inches(px_value, axis_size_px):
    """像素转英寸的精确转换"""
    if axis_size_px == CANVAS_WIDTH_PX:
        return Inches(px_value * SLIDE_WIDTH_INCHES / CANVAS_WIDTH_PX)
    else:
        return Inches(px_value * SLIDE_HEIGHT_INCHES / CANVAS_HEIGHT_PX)

def rgb_string_to_pptx_color(rgb_str):
    """RGB字符串转RGBColor对象"""
    try:
        r, g, b = map(int, rgb_str.strip('RGB()').split(','))
        return RGBColor(r, g, b)
    except:
        return RGBColor(0, 0, 0)

def get_shape_center(shape):
    """获取形状中心点坐标 (EMU)"""
    center_x = shape.left + shape.width / 2
    center_y = shape.top + shape.height / 2
    return center_x, center_y

def get_shape_bounds(shape):
    """获取形状边界信息 (EMU)"""
    return {
        'left': shape.left,
        'right': shape.left + shape.width,
        'top': shape.top,
        'bottom': shape.top + shape.height,
        'center_x': shape.left + shape.width / 2,
        'center_y': shape.top + shape.height / 2,
        'width': shape.width,
        'height': shape.height
    }

def calculate_distance(start_shape, end_shape):
    """计算两个形状中心点之间的距离 (EMU)"""
    start_x, start_y = get_shape_center(start_shape)
    end_x, end_y = get_shape_center(end_shape)
    return math.sqrt((end_x - start_x)**2 + (end_y - start_y)**2)

def get_connection_point_coords(bounds, connection_point):
    """根据连接点索引获取实际坐标 (EMU)"""
    if connection_point == CONNECT_LEFT:
        return (bounds['left'], bounds['center_y'])
    elif connection_point == CONNECT_TOP:
        return (bounds['center_x'], bounds['top'])
    elif connection_point == CONNECT_RIGHT:
        return (bounds['right'], bounds['center_y'])
    elif connection_point == CONNECT_BOTTOM:
        return (bounds['center_x'], bounds['bottom'])
    else:
        return (bounds['center_x'], bounds['center_y'])


def calculate_connection_score(start_bounds, end_bounds, start_conn, end_conn):
    """
    计算连接方案的美观度分数（越小越好）。
    """
    start_point = get_connection_point_coords(start_bounds, start_conn)
    end_point = get_connection_point_coords(end_bounds, end_conn)
    
    dx_len = end_point[0] - start_point[0]
    dy_len = end_point[1] - start_point[1]

    line_length = math.sqrt(dx_len**2 + dy_len**2)
    
    try:
        SMALL_TOLERANCE_EMU = Inches(0.1)._emu
    except:
        SMALL_TOLERANCE_EMU = 1000

    angle_penalty = 0
    if abs(dx_len) > SMALL_TOLERANCE_EMU and abs(dy_len) > SMALL_TOLERANCE_EMU:
        # 对角线有轻微惩罚
        angle_penalty = min(abs(dx_len), abs(dy_len)) * 0.3
    
    direction_bonus = 0
    center_dx = end_bounds['center_x'] - start_bounds['center_x']
    center_dy = end_bounds['center_y'] - start_bounds['center_y']
    
    # 方向一致性奖励
    if ((start_conn == CONNECT_RIGHT and end_conn == CONNECT_LEFT and center_dx > 0) or
        (start_conn == CONNECT_LEFT and end_conn == CONNECT_RIGHT and center_dx < 0) or
        (start_conn == CONNECT_BOTTOM and end_conn == CONNECT_TOP and center_dy > 0) or
        (start_conn == CONNECT_TOP and end_conn == CONNECT_BOTTOM and center_dy < 0)):
        direction_bonus = -line_length * 0.2
    
    return line_length + angle_penalty + direction_bonus


def analyze_element_relationship(start_shape, end_shape, start_id, end_id):
    """详细分析两个元素的位置关系，并返回推荐连接点"""
    start_bounds = get_shape_bounds(start_shape)
    end_bounds = get_shape_bounds(end_shape)
    
    X_TOLERANCE = get_emu_value(X_ALIGNMENT_TOLERANCE_IN)
    Y_TOLERANCE = get_emu_value(Y_ALIGNMENT_TOLERANCE_IN)
    MIN_SEPARATION_EMU = get_emu_value(Inches(0.1)) 
    
    dx = end_bounds['center_x'] - start_bounds['center_x']
    dy = end_bounds['center_y'] - start_bounds['center_y']
    abs_dx = abs(dx)
    abs_dy = abs(dy)
    euclidean_distance = math.sqrt(dx**2 + dy**2)
    x_overlap = (start_bounds['right'] > end_bounds['left'] and start_bounds['left'] < end_bounds['right'])
    y_overlap = (start_bounds['bottom'] > end_bounds['top'] and start_bounds['top'] < end_bounds['bottom'])
    direction_ratio = abs_dy / abs_dx if abs_dx > 1 else float('inf') 
    
    EMU_PER_INCH = 914400 
    print(f"\n🔍 分析连接: {start_id} -> {end_id}")
    print(f"    中心距离: dx={dx/EMU_PER_INCH:.2f}英寸, dy={dy/EMU_PER_INCH:.2f}英寸")
    
    # --- 分析逻辑（精简输出，与之前逻辑一致） ---
    recommended_connection = (CONNECT_RIGHT, CONNECT_LEFT) # 默认值
    relationship_type = "对角-均衡"

    if x_overlap and abs_dy > Y_TOLERANCE:
        if dy > 0: recommended_connection = (CONNECT_BOTTOM, CONNECT_TOP)
        else: recommended_connection = (CONNECT_TOP, CONNECT_BOTTOM)
    elif abs_dx < X_TOLERANCE:
        if dy > 0: recommended_connection = (CONNECT_BOTTOM, CONNECT_TOP)
        else: recommended_connection = (CONNECT_TOP, CONNECT_BOTTOM)
    elif y_overlap and abs_dx > X_TOLERANCE:
        if dx > 0: recommended_connection = (CONNECT_RIGHT, CONNECT_LEFT)
        else: recommended_connection = (CONNECT_LEFT, CONNECT_RIGHT)
    elif abs_dy < Y_TOLERANCE:
        if dx > 0: recommended_connection = (CONNECT_RIGHT, CONNECT_LEFT)
        else: recommended_connection = (CONNECT_LEFT, CONNECT_RIGHT)
    elif direction_ratio > 1.5:
        if dy > 0: recommended_connection = (CONNECT_BOTTOM, CONNECT_TOP)
        else: recommended_connection = (CONNECT_TOP, CONNECT_BOTTOM)
    elif direction_ratio < 0.67:
        if dx > 0: recommended_connection = (CONNECT_RIGHT, CONNECT_LEFT)
        else: recommended_connection = (CONNECT_LEFT, CONNECT_RIGHT)
    else:
        # 对角关系倾向于最远的轴
        if abs_dy > abs_dx:
            if dy > 0: recommended_connection = (CONNECT_BOTTOM, CONNECT_TOP)
            else: recommended_connection = (CONNECT_TOP, CONNECT_BOTTOM)
        else:
            if dx > 0: recommended_connection = (CONNECT_RIGHT, CONNECT_LEFT)
            else: recommended_connection = (CONNECT_LEFT, CONNECT_RIGHT)
    
    connection_names = {CONNECT_LEFT: "左侧", CONNECT_TOP: "顶部", CONNECT_RIGHT: "右侧", CONNECT_BOTTOM: "底部"}
    start_conn_name = connection_names[recommended_connection[0]]
    end_conn_name = connection_names[recommended_connection[1]]
    print(f"    ✅ 推荐连接: {start_id}的{start_conn_name} -> {end_id}的{end_conn_name}")
    
    return relationship_type, recommended_connection

def get_smart_connection_points_enhanced(start_shape, end_shape, start_id="", end_id=""):
    """
    增强版智能连接点选择：进行详细分析并根据分数选择最优连接点。
    """
    print(f"\n🧠 使用增强版智能连接分析...")
    
    # 1. 首先使用基础分析获得推荐连接
    relationship_type, base_recommendation = analyze_element_relationship(
        start_shape, end_shape, start_id, end_id
    )
    
    start_bounds = get_shape_bounds(start_shape)
    end_bounds = get_shape_bounds(end_shape)
    
    # 2. 计算所有可能连接点组合的"美观度分数"
    connection_options = [
        (CONNECT_LEFT, CONNECT_RIGHT), (CONNECT_RIGHT, CONNECT_LEFT),
        (CONNECT_TOP, CONNECT_BOTTOM), (CONNECT_BOTTOM, CONNECT_TOP),
        # 考虑对角连接
        (CONNECT_RIGHT, CONNECT_TOP), (CONNECT_RIGHT, CONNECT_BOTTOM),
        (CONNECT_LEFT, CONNECT_TOP), (CONNECT_LEFT, CONNECT_BOTTOM),
        (CONNECT_TOP, CONNECT_LEFT), (CONNECT_TOP, CONNECT_RIGHT),
        (CONNECT_BOTTOM, CONNECT_LEFT), (CONNECT_BOTTOM, CONNECT_RIGHT),
    ]
    
    best_score = float('inf')
    best_connection = base_recommendation
    scores_info = []
    
    for (start_conn, end_conn) in connection_options:
        score = calculate_connection_score(start_bounds, end_bounds, start_conn, end_conn)
        
        # 如果是基础分析推荐的连接，给予额外奖励
        if (start_conn, end_conn) == base_recommendation:
            score *= 0.8
            
        scores_info.append((start_conn, end_conn, score))
        
        if score < best_score:
            best_score = score
            best_connection = (start_conn, end_conn)
    
    # 输出评分详情
    scores_info.sort(key=lambda x: x[2])
    connection_names = {
        CONNECT_LEFT: "左", CONNECT_TOP: "上", 
        CONNECT_RIGHT: "右", CONNECT_BOTTOM: "下"
    }
    
    try:
        EMU_PER_INCH = 914400 
    except:
        EMU_PER_INCH = 1
        
    print(f"    🏆 最优连接方案 (前5名):")
    for i, (start_conn, end_conn, score) in enumerate(scores_info[:5]):
        start_name = connection_names[start_conn]
        end_name = connection_names[end_conn]
        is_chosen = "✅" if (start_conn, end_conn) == best_connection else "  "
        print(f"    {is_chosen} {i+1}. {start_name}→{end_name}: {score/EMU_PER_INCH:.2f} (近似英寸)")
    
    return best_connection


def load_diagram_data():
    """安全加载JSON数据"""
    try:
        # 假设 data 文件名为 diagram_data.json
        with open('diagram_data.json', 'r', encoding='utf-8') as f: 
            return json.load(f)
    except FileNotFoundError:
        print("错误：未找到 'diagram_data.json' 文件。请确保该文件存在于脚本的同一目录下。")
        return None
    except json.JSONDecodeError as e:
        print(f"错误：JSON文件格式不正确 - {e}")
        return None


# ----------------------------------------------------
# --- 核心函数定义 ---
# ----------------------------------------------------

def create_powerpoint_diagram(data):
    """根据JSON数据创建PowerPoint图表，并使用强制坐标设置连接线。"""
    if not data:
        return False
    
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_WIDTH_INCHES)
    prs.slide_height = Inches(SLIDE_HEIGHT_INCHES)
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    shapes_map = {}
    
    # 第一步：创建所有形状元素
    print("正在创建形状元素...")
    for element in data['elements']:
        if not all(key in element for key in ['text', 'position', 'dimensions']):
            continue
            
        try:
            x_px_str, y_px_str = element['position'].strip('[]').split(',')
            w_px_str, h_px_str = element['dimensions'].strip('[]').split(',')
            x_px, y_px = int(x_px_str.strip()), int(y_px_str.strip())
            w_px, h_px = int(w_px_str.strip()), int(h_px_str.strip())
            
            # 转换为英寸
            left = px_to_inches(x_px, CANVAS_WIDTH_PX)
            top = px_to_inches(y_px, CANVAS_HEIGHT_PX)
            width = px_to_inches(w_px, CANVAS_WIDTH_PX)
            height = px_to_inches(h_px, CANVAS_HEIGHT_PX) 
            
        except (ValueError, IndexError) as e:
            print(f"解析元素位置失败: {element.get('text', 'Unknown')} - {e}")
            continue

        shape_type = MSO_AUTO_SHAPE_TYPE.RECTANGLE
        element_type = element.get('type', '')
        
        if 'rounded_rectangle' in element_type:
            shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE
        elif 'circle' in element_type or 'oval' in element_type:
            shape_type = MSO_AUTO_SHAPE_TYPE.OVAL
        
        if 'text' in element_type and 'independent_text' in element_type:
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.text = element['text']
            tf.paragraphs[0].font.size = Pt(9)
            continue
        
        shape = slide.shapes.add_shape(shape_type, left, top, width, height)

        fill = shape.fill
        fill.solid()
        fill.fore_color.rgb = rgb_string_to_pptx_color(element.get('color', 'RGB(200, 200, 200)'))
        line = shape.line
        line.color.rgb = RGBColor(0, 0, 0)
        line.width = Pt(0.75)

        text_frame = shape.text_frame
        text_frame.text = element['text']
        text_frame.word_wrap = True
        text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        paragraph = text_frame.paragraphs[0]
        paragraph.font.size = Pt(10)
        paragraph.font.color.rgb = RGBColor(0, 0, 0)
        paragraph.alignment = MSO_ANCHOR.TOP
        
        if 'id' in element:
            shapes_map[element['id']] = shape
            
    print(f"成功创建 {len(shapes_map)} 个形状元素")
    
    # 第二步：创建连接线 (使用强制坐标设置)
    print("正在创建连接线...")
    connection_count = 0
    min_dist_emu = get_emu_value(MIN_DISTANCE_THRESHOLD_IN)
    
    for relationship in data.get('relationships', []):
        from_id = relationship.get('from')
        to_id = relationship.get('to')
        link_type = relationship.get('type', '')
        
        if from_id not in shapes_map or to_id not in shapes_map:
            if from_id and to_id:
                print(f"警告：找不到连接的形状 {from_id} -> {to_id}")
            continue
            
        start_shape = shapes_map[from_id]
        end_shape = shapes_map[to_id]
        
        if calculate_distance(start_shape, end_shape) < min_dist_emu:
            print(f"跳过距离过近的连接: {from_id} -> {to_id}")
            continue

        if link_type in ['arrow', 'line', 'arrow_flow_down', 'flow'] or 'loose_line' in link_type:
            
            # 1. 获取分析出的最优连接点
            start_conn_pt, end_conn_pt = get_smart_connection_points_enhanced(
                start_shape, end_shape, from_id, to_id
            )
            
            # 2. 获取连接点的实际坐标 (EMU)
            start_bounds = get_shape_bounds(start_shape)
            end_bounds = get_shape_bounds(end_shape)
            start_coord = get_connection_point_coords(start_bounds, start_conn_pt)
            end_coord = get_connection_point_coords(end_bounds, end_conn_pt)

            try:
                # 3. 创建连接器 (使用直线类型，方便手动设置坐标)
                connector = slide.shapes.add_connector(
                    CONNECTOR_STRAIGHT, Inches(0), Inches(0), Inches(0), Inches(0)
                )

                # 4. **必须：建立形状间的逻辑连接**
                #    这一步保证了形状被拖动时，连接线也能自动跟随。
                connector.begin_connect(start_shape, start_conn_pt)
                connector.end_connect(end_shape, end_conn_pt)
                
                # 5. **核心修正：强制设置连接器的几何坐标**
                #    这一步确保连接线在 PPT 中显示时，其端点位置是精确计算的。
                
                left_emu = min(start_coord[0], end_coord[0])
                top_emu = min(start_coord[1], end_coord[1])
                width_emu = abs(start_coord[0] - end_coord[0])
                height_emu = abs(start_coord[1] - end_coord[1])
                
                # 直接修改形状元素（sp.x/y/cx/cy 对应 left/top/width/height）
                sp = connector.element
                sp.x, sp.y = int(left_emu), int(top_emu)
                sp.cx, sp.cy = int(width_emu), int(height_emu)
                
                # 6. 设置线条样式
                line = connector.line
                line.color.rgb = RGBColor(0, 0, 0)
                line.width = Pt(1.5)
                
                is_loose = 'loose_line' in link_type
                
                # 设置箭头
                # 修正：通常箭头线是主要流程，无箭头线是辅助或松散连接。
                if 'arrow' in link_type or ('flow' in link_type and not is_loose):
                    line.end_arrowhead_style = ARROWHEAD_TRIANGLE
                else:
                    line.end_arrowhead_style = ARROWHEAD_NONE
                
                # 松散连接的特殊样式
                if is_loose:
                    line.color.rgb = RGBColor(128, 128, 128)
                    line.width = Pt(1)
                    print(f"    ✅ 松散连接创建成功 (强制坐标)")
                else:
                    print(f"    ✅ 连接创建成功 (强制坐标)")
                    
                connection_count += 1
                
            except Exception as e:
                print(f"    ❌ 连接失败: {e}")

    print(f"成功创建 {connection_count} 个连接")
    
    # 保存文件
    output_path = '优化版_实例映射研讨会结果.pptx'
    try:
        prs.save(output_path)
        print(f"\n✅ 成功生成PowerPoint文件: {output_path}")
        return True
    except Exception as e:
        print(f"❌ 保存文件失败: {e}")
        return False


def main():
    """主函数"""
    print("🚀 开始生成PowerPoint图表...")
    print("=" * 60)
    
    data = load_diagram_data()
    if not data:
        return
    
    print(f"📊 数据加载成功:")
    print(f"    - 元素数量: {len(data.get('elements', []))}")
    print(f"    - 关系数量: {len(data.get('relationships', []))}")
    print("=" * 60)
    
    success = create_powerpoint_diagram(data)
    
    print("=" * 60)
    if success:
        print("🎉 图表生成完成！")
    else:
        print("❌ 图表生成失败，请检查上述错误信息。")

if __name__ == "__main__":
    main()